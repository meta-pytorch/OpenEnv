from __future__ import annotations

import re
from pathlib import Path
from typing import Any


def strip_html(text: str) -> str:
    """Turn simple HTML used in the web deck into plain text for PowerPoint."""
    t = text.replace("<br>", "\n").replace("<br/>", "\n")
    t = re.sub(r"<p[^>]*>", "", t, flags=re.I)
    t = t.replace("</p>", "\n\n")
    t = re.sub(r"<strong>(.*?)</strong>", r"\1", t, flags=re.I | re.S)
    t = re.sub(r"<em>(.*?)</em>", r"\1", t, flags=re.I | re.S)
    t = re.sub(r"<code>(.*?)</code>", r"\1", t, flags=re.I | re.S)
    t = re.sub(r"<[^>]+>", "", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


# Simple advisor deck: fewer slides, bigger ideas, minimal jargon.
SLIDES: list[dict[str, Any]] = [
    {
        "kind": "cover",
        "title": "Pathway analysis in OpenEnv",
        "subtitle": "A simple, repeatable framework to run RNA-seq → DE → pathways (and to evaluate AI copilots fairly).",
        "meta": "Parul Sarma Aryasomayajula — advisor meeting",
    },
    {
        "kind": "content",
        "section": "Problem",
        "title": "Why this matters",
        "bullets": [
            "Biology conclusions often depend on analysis choices (contrast, filtering, gene sets).",
            "AI copilots are coming, but we need a reliable way to check if they are correct.",
            "Right now: results can be hard to reproduce and hard to compare across tools.",
        ],
    },
    {
        "kind": "content",
        "section": "Idea",
        "title": "What OpenEnv gives us (in plain terms)",
        "bullets": [
            "A standard “container” for an analysis workflow: same inputs, same steps, same outputs.",
            "A record of what happened (like a computational lab notebook).",
            "A way to compare two copilots on the exact same task, fairly.",
        ],
    },
    {
        "kind": "content",
        "section": "My build",
        "title": "What I built",
        "bullets": [
            "A pathway analysis environment: counts + sample labels → differential expression → pathway enrichment.",
            "A simple UI (Gradio) to demo the workflow quickly.",
            "Saved outputs (tables + JSON summaries) so results can be reviewed and rerun.",
        ],
    },
    {
        "kind": "content",
        "section": "Demo flow",
        "title": "What happens in one run (one study)",
        "bullets": [
            "Choose a study + define the comparison (e.g., treated vs control).",
            "Run differential expression to get a ranked gene list.",
            "Run pathway enrichment to get the main biological themes.",
            "Export a short report + artifacts for review.",
        ],
    },
    {
        "kind": "table",
        "section": "Why not a script?",
        "title": "Script vs OpenEnv-style workflow",
        "headers": ("A normal script", "This OpenEnv environment"),
        "rows": [
            ("Runs once, hard to compare", "Runs many times, comparable across models"),
            ("Hard to audit what happened", "Trace + saved artifacts for audit"),
            ("Hard to reuse as a benchmark", "Designed to be a benchmark"),
        ],
    },
    {
        "kind": "content",
        "section": "Evidence",
        "title": "How I ground it in real public data (GEO)",
        "bullets": [
            "I used multiple GEO-style study types, because public data is messy.",
            "Best case: real integer counts → strongest validation.",
            "Other cases: author DE tables or TPM-only supplements → useful for testing, labeled honestly.",
        ],
    },
    {
        "kind": "content",
        "section": "Limitations",
        "title": "What I will NOT over-claim",
        "bullets": [
            "Small sample sizes mean noisy DE: pathways are hypothesis-level, not final truth.",
            "Some GEO studies do not provide ideal raw counts; I label these cases clearly.",
            "Heavy raw-read workflows (SRA → quant) take real compute time.",
        ],
    },
    {
        "kind": "content",
        "section": "Ask / next",
        "title": "What I want feedback on",
        "bullets": [
            "What is the best advisor narrative: evaluation benchmark, teaching tool, or methods paper?",
            "Which biological study archetypes would be most convincing to add next?",
            "What should be the simplest success metric for a first paper/demo?",
        ],
    },
    {
        "kind": "cover",
        "title": "Discussion",
        "subtitle": "Happy to walk through a single GEO study end-to-end and discuss what to claim (and what not to claim).",
        "meta": "Thank you",
    },
]


def _fill_bullets(text_frame: Any, lines: list[str]) -> None:
    text_frame.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.level = 0


def build_pptx(out_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover
        raise SystemExit(
            "Missing dependency python-pptx. Install with: uv pip install python-pptx"
        ) from exc

    prs = Presentation()
    prs.slide_width = 13_333_333
    prs.slide_height = 7_500_000

    layout_title = prs.slide_layouts[0]
    layout_content = prs.slide_layouts[1]
    layout_two = prs.slide_layouts[3]
    layout_title_only = prs.slide_layouts[5]

    for spec in SLIDES:
        kind = spec["kind"]
        if kind == "cover":
            slide = prs.slides.add_slide(layout_title)
            slide.shapes.title.text = spec["title"]
            sub = slide.placeholders[1]
            sub.text = spec["subtitle"] + "\n\n" + spec.get("meta", "")
            continue

        if kind == "content":
            slide = prs.slides.add_slide(layout_content)
            sec = spec.get("section", "")
            slide.shapes.title.text = f"{sec}: {spec['title']}" if sec else spec["title"]
            body = slide.shapes.placeholders[1].text_frame
            bullets = [strip_html(b) for b in spec["bullets"]]
            _fill_bullets(body, bullets)
            continue

        if kind == "two_col":
            slide = prs.slides.add_slide(layout_two)
            sec = spec.get("section", "")
            slide.shapes.title.text = f"{sec}: {spec['title']}" if sec else spec["title"]
            left_tf = slide.placeholders[1].text_frame
            right_tf = slide.placeholders[2].text_frame
            left_lines = [spec["left_title"]] + [strip_html(x) for x in spec["left"]]
            right_lines = [spec["right_title"]] + [strip_html(x) for x in spec["right"]]
            _fill_bullets(left_tf, left_lines)
            _fill_bullets(right_tf, right_lines)
            for i, para in enumerate(left_tf.paragraphs):
                para.level = 0 if i == 0 else 1
                if i == 0:
                    para.font.bold = True
            for i, para in enumerate(right_tf.paragraphs):
                para.level = 0 if i == 0 else 1
                if i == 0:
                    para.font.bold = True
            continue

        if kind == "table":
            slide = prs.slides.add_slide(layout_title_only)
            slide.shapes.title.text = f"{spec['section']}: {spec['title']}"
            headers = spec["headers"]
            rows_data: list[tuple[str, str]] = [headers, *spec["rows"]]
            nrows, ncols = len(rows_data), 2
            left, top, width, height = Inches(0.6), Inches(1.35), Inches(12.2), Inches(5.2)
            table_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
            tbl = table_shape.table
            for r, row in enumerate(rows_data):
                for c, cell in enumerate(row):
                    tbl.cell(r, c).text = strip_html(str(cell))
                    for paragraph in tbl.cell(r, c).text_frame.paragraphs:
                        paragraph.font.size = Pt(14 if r > 0 else 15)
                        if r == 0:
                            paragraph.font.bold = True
            continue

        if kind == "diagram":
            slide = prs.slides.add_slide(layout_content)
            sec = spec.get("section", "")
            slide.shapes.title.text = f"{sec}: {spec['title']}"
            body = slide.shapes.placeholders[1].text_frame
            _fill_bullets(body, [strip_html(b) for b in spec["bullets"]])
            continue

        if kind == "quote":
            slide = prs.slides.add_slide(layout_content)
            sec = spec.get("section", "")
            slide.shapes.title.text = f"{sec}: {spec['title']}"
            body = slide.shapes.placeholders[1].text_frame
            plain = strip_html(spec["html"])
            parts = [p.strip() for p in plain.split("\n\n") if p.strip()]
            _fill_bullets(body, parts)
            for para in body.paragraphs:
                para.font.size = Pt(15)
            continue

        raise ValueError(f"Unknown slide kind: {kind}")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    env_root = Path(__file__).resolve().parents[1]
    out = env_root / "docs" / "ADVISOR_PRESENTATION_SIMPLE.pptx"
    build_pptx(out)
    print(str(out))


if __name__ == "__main__":
    main()
