from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class Slide:
    title: str
    bullets: list[str]
    subtitle: str | None = None


SLIDES: list[Slide] = [
    Slide(
        title="OpenEnv: Pathway Analysis Environment",
        subtitle="A realistic RNA-seq workflow environment for evaluating scientific agents",
        bullets=["Parul Sarma Aryasomayajula"],
    ),
    Slide(
        title="The problem",
        bullets=[
            "Scientific/biomed agent performance is hard to evaluate",
            "Real workflows require multi-step tool use, not single-shot answers",
            "We need reproducible, instrumented tasks with measurable outcomes",
            "Goal: evaluate agents on end-to-end reasoning + tool execution",
        ],
    ),
    Slide(
        title="What OpenEnv provides",
        bullets=[
            "Standard environment API (reset / step / state)",
            "Client/server separation (FastAPI runtime; reusable clients)",
            "Structured observations & metadata for evaluation and debugging",
            "Makes complex tasks testable like environments",
        ],
    ),
    Slide(
        title="Task: pathway analysis from gene expression",
        bullets=[
            "Objective: infer what biological programs/pathways explain a phenotype",
            "Workflow: design → DGE → pathway enrichment → compare → hypothesis",
        ],
    ),
    Slide(
        title="Agent interface (actions)",
        bullets=[
            "understand_experiment_design",
            "inspect_dataset",
            "run_differential_expression",
            "run_pathway_enrichment",
            "compare_pathways",
            "ask_expert (budgeted hint/curriculum)",
            "submit_answer",
        ],
    ),
    Slide(
        title="Observability & safety rails",
        bullets=[
            "DE table: log2FC, p-value, q-value, baseMean",
            "Enriched pathways + overlap genes",
            "Ambiguity/overlap summaries when multiple hits look similar",
            "Stable failure codes for error analysis",
            "HTML episode trace for debugging & review",
        ],
    ),
    Slide(
        title="Real pipeline integrated (not a toy)",
        bullets=[
            "Differential expression: DESeq2-style stats via PyDESeq2",
            "DESeq2 prefiltering for low-count genes",
            "Configurable thresholds (FDR, min log2FC, direction)",
            "Pathway analysis: ORA on query genes",
            "Supports case-defined gene sets (offline) or Enrichr via gseapy",
        ],
    ),
    Slide(
        title="Real-data support (GEO-style inputs)",
        bullets=[
            "Counts files like *.csv.gz (GEO)",
            "Sample metadata (condition labels)",
            "Pipeline cases defined as JSON configs",
            "Same agent actions run on real public studies",
        ],
    ),
    Slide(
        title="Public study demo: GSE235417",
        bullets=[
            "HNSCC PDX (UCLHN04): baseline vs acquired cetuximab-resistant",
            "3 biological replicates per condition",
            "DESeq2 contrast: resistant vs baseline",
            "Enrichment: Hallmark / KEGG / Reactome",
            "Exports: JSON results + HTML trace + report",
        ],
    ),
    Slide(
        title="Example results (headlines)",
        bullets=[
            "Top DE genes (example): DAPK1, CXCL8, SOX2, FKBP5, …",
            "Top pathways (example): EMT; inflammatory/interferon; TNFα/NF-κB; cytokine signaling",
            "Emphasis: reproducible pipeline + measurable signals (avoid over-claiming biology)",
        ],
    ),
    Slide(
        title="What we can measure (evaluation metrics)",
        bullets=[
            "Correctness: does the agent identify the right hypothesis?",
            "Efficiency: steps taken, tool calls, reward trajectory",
            "Robustness: invalid contrasts, missing metadata, low signal",
            "Interpretability: traces + structured evidence for postmortems",
        ],
    ),
    Slide(
        title="Why this is valuable",
        bullets=[
            "Bridges toy benchmarks and real science workflows",
            "Makes agent behavior reproducible, debuggable, comparable",
            "Foundation for regression tests, eval harnesses, future RL-style loops",
        ],
    ),
    Slide(
        title="Roadmap",
        bullets=[
            "Add GSEA (rank-based enrichment) alongside ORA",
            "Richer experiment designs (covariates / batches)",
            "More datasets and tasks (perturbations, multi-omics, drug response)",
            "Standardized scoring + reporting for agent comparisons",
        ],
    ),
    Slide(
        title="Close",
        bullets=[
            "OpenEnv environments turn real workflows into measurable evaluation problems",
            "Happy to discuss fit with your agent evaluation stack and next domains",
        ],
    ),
]


def build_pptx(out_path: Path) -> None:
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover
        raise SystemExit(
            "Missing dependency python-pptx. Install with: uv pip install python-pptx"
        ) from exc

    prs = Presentation()
    # 16:9 wide
    prs.slide_width = 13_333_333
    prs.slide_height = 7_500_000

    title_and_content = prs.slide_layouts[1]

    for s in SLIDES:
        slide = prs.slides.add_slide(title_and_content)
        slide.shapes.title.text = s.title
        body = slide.shapes.placeholders[1].text_frame
        body.clear()

        if s.subtitle:
            p = body.paragraphs[0]
            p.text = s.subtitle
            p.level = 0
            for b in s.bullets:
                q = body.add_paragraph()
                q.text = b
                q.level = 0
        else:
            for i, b in enumerate(s.bullets):
                if i == 0:
                    p = body.paragraphs[0]
                    p.text = b
                    p.level = 0
                else:
                    q = body.add_paragraph()
                    q.text = b
                    q.level = 0

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    repo_root = Path(__file__).resolve().parents[2]
    out = repo_root / "docs" / "INTERVIEW_PITCH_DECK.pptx"
    build_pptx(out)
    print(str(out))


if __name__ == "__main__":
    main()

